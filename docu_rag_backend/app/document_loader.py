import os
from langchain_core.documents import Document
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from PIL import Image
import pytesseract
import whisper
import langid
import pytesseract
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

def detect_language(text):
    lang, _ = langid.classify(text)
    return lang


def load_document(file_path):
    ext = file_path.lower().split(".")[-1]
    filename = os.path.basename(file_path)

    # -------- PDF --------
    if ext == "pdf":
        reader = PdfReader(file_path)
        text = ""
        for page in reader.pages:
            text += page.extract_text() or ""

    # -------- DOCX --------
    elif ext == "docx":
        doc = DocxDocument(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])

    # -------- TXT --------
    elif ext == "txt":
        with open(file_path, "r", encoding="utf-8") as f:
            text = f.read()

    # -------- PPTX --------
    elif ext == "pptx":
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

    # -------- IMAGE --------
    elif ext in ["png", "jpg", "jpeg"]:
        image = Image.open(file_path)
        text = pytesseract.image_to_string(image)

    # -------- AUDIO --------
    elif ext in ["mp3", "wav"]:
        model = whisper.load_model("base")
        result = model.transcribe(file_path)
        text = result["text"]

    else:
        raise ValueError("Unsupported file type")

    lang = detect_language(text)

    return [
        Document(
            page_content=text,
            metadata={
                "source": filename,
                "language": lang
            }
        )
    ]