# =============================================================================
#  Gyana AI  –  Document Service
#  Extracts clean text from PDF, DOCX, PPTX, TXT
#  PDF:  PyMuPDF page-by-page (memory efficient for large files)
#  DOCX: python-docx – paragraphs + headings + tables
#  PPTX: python-pptx – all shapes + tables + speaker notes
#  TXT:  UTF-8 with latin-1 fallback
# =============================================================================

from __future__ import annotations

import logging
import re
from pathlib import Path

log = logging.getLogger("gyana.document")

# Max pages to process (prevents hanging on huge PDFs)
MAX_PAGES = 100


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------
def extract_text(file_path: str | Path) -> str:
    path = Path(file_path)
    ext  = path.suffix.lower()

    extractors = {
        ".pdf":  _extract_pdf,
        ".docx": _extract_docx,
        ".pptx": _extract_pptx,
        ".txt":  _extract_txt,
    }

    fn = extractors.get(ext)
    if fn is None:
        raise ValueError(f"No text extractor for extension '{ext}'.")

    raw = fn(path)
    return _clean(raw)


# ---------------------------------------------------------------------------
# PDF — page by page with PyMuPDF
# ---------------------------------------------------------------------------
def _extract_pdf(path: Path) -> str:
    try:
        import fitz  # PyMuPDF

        parts = []
        with fitz.open(str(path)) as doc:
            total_pages = len(doc)
            pages_to_process = min(total_pages, MAX_PAGES)

            if total_pages > MAX_PAGES:
                log.warning("PDF has %d pages, processing first %d only", total_pages, MAX_PAGES)

            for i in range(pages_to_process):
                try:
                    page = doc[i]
                    text = page.get_text("text")
                    if text.strip():
                        parts.append(f"[Page {i+1}]\n{text}")
                    # Free page memory immediately
                    page = None
                except Exception as exc:
                    log.warning("Skipping page %d: %s", i+1, exc)
                    continue

        result = "\n\n".join(parts)
        log.info("PDF extracted: %d chars from %d pages", len(result), pages_to_process)
        return result

    except ImportError:
        raise RuntimeError("pip install pymupdf")
    except Exception as exc:
        raise RuntimeError(f"PDF extraction failed: {exc}")


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------
def _extract_docx(path: Path) -> str:
    try:
        from docx import Document
    except ImportError:
        raise RuntimeError("pip install python-docx")

    doc   = Document(str(path))
    parts: list[str] = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = para.style.name if para.style else ""
        if style.startswith("Heading"):
            lvl = "".join(filter(str.isdigit, style)) or "1"
            parts.append(f"{'#' * int(lvl)} {text}")
        else:
            parts.append(text)

    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(
                cell.text.strip() for cell in row.cells if cell.text.strip()
            )
            if row_text:
                parts.append(row_text)

    return "\n\n".join(parts)


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------
def _extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("pip install python-pptx")

    prs   = Presentation(str(path))
    parts: list[str] = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_parts = [f"[Slide {slide_num}]"]

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide_parts.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(
                        c.text.strip() for c in row.cells if c.text.strip()
                    )
                    if row_text:
                        slide_parts.append(row_text)

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_parts.append(f"[Notes] {notes}")

        if len(slide_parts) > 1:
            parts.append("\n".join(slide_parts))

    return "\n\n".join(parts)


# ---------------------------------------------------------------------------
# TXT
# ---------------------------------------------------------------------------
def _extract_txt(path: Path) -> str:
    for enc in ("utf-8", "utf-8-sig", "latin-1", "cp1252"):
        try:
            return path.read_text(encoding=enc)
        except (UnicodeDecodeError, LookupError):
            continue
    raise RuntimeError(f"Cannot decode text file: {path.name}")


# ---------------------------------------------------------------------------
# Text cleaning
# ---------------------------------------------------------------------------
def _clean(text: str) -> str:
    if not text:
        return ""
    text = text.replace("\x00", "")
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"\n{3,}", "\n\n", text)
    lines = [ln.rstrip() for ln in text.split("\n")]
    return "\n".join(lines).strip()